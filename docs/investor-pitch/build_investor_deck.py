#!/usr/bin/env python3
"""
Generate Growl investor pitch deck (PowerPoint).

Usage (from repo root):
  pip install python-pptx
  python3 docs/investor-pitch/build_investor_deck.py

Output: docs/investor-pitch/Growl_Investor_Deck.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT_PATH = Path(__file__).resolve().parent / "Growl_Investor_Deck.pptx"

ACCENT = RGBColor(0x05, 0x96, 0x69)  # emerald-600
MUTED = RGBColor(0x57, 0x53, 0x4E)


def _title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    box = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(8.5), Inches(1.4))
    tf = box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.LEFT

    sub = slide.shapes.add_textbox(Inches(0.7), Inches(3.6), Inches(8.5), Inches(1.2))
    stf = sub.text_frame
    stf.text = subtitle
    sp = stf.paragraphs[0]
    sp.font.size = Pt(18)
    sp.font.color.rgb = MUTED


def _bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    t = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(8.5), Inches(0.8))
    tf = t.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT

    body = slide.shapes.add_textbox(Inches(0.85), Inches(1.35), Inches(8.2), Inches(5.5))
    btf = body.text_frame
    btf.word_wrap = True
    for i, line in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(17)
        p.font.color.rgb = RGBColor(0x29, 0x25, 0x24)
        p.space_after = Pt(10)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    _title_slide(
        prs,
        "Growl",
        "Social growth, marketplace, and business tools — one mobile platform.\nInvestor overview · 2026",
    )

    _bullet_slide(
        prs,
        "The opportunity",
        [
            "Consumers want habit and skill growth with community accountability — not only passive content.",
            "Creators and local businesses need lightweight commerce, not another full e-commerce build.",
            "A unified feed + marketplace + business dashboard reduces fragmentation and increases retention.",
        ],
    )

    _bullet_slide(
        prs,
        "What Growl is",
        [
            "Mobile-first social feed for progress posts, reactions, and discovery by category.",
            "Stories and short-form surfaces for engagement (Reels-style experience on the roadmap).",
            "Marketplace for products with orders and line items tied to business accounts.",
            "Business dashboard: KPIs, inventory, order status — API-backed and shipping today.",
        ],
    )

    _bullet_slide(
        prs,
        "Product status",
        [
            "Backend: Cloudflare Workers + D1 + KV; REST API at /api/v1 with automated tests.",
            "Auth: email/password and SSO hooks; role-aware flows for business and instructors.",
            "Frontend: React Native (Expo SDK 54), Redux Toolkit for global state, polished feed and marketing UX.",
            "Health, feed, marketplace, and business routes verified against production-like environments.",
        ],
    )

    _bullet_slide(
        prs,
        "Why now",
        [
            "Serverless edge APIs reduce cost to serve and speed iteration for early-stage teams.",
            "Cross-platform (iOS / Android / web-capable stack) maximizes addressable audience.",
            "Community + commerce convergence is where modern consumer apps capture LTV.",
        ],
    )

    _bullet_slide(
        prs,
        "Business model (directional)",
        [
            "Take rate on marketplace transactions and sponsored placements in feed/marketing surfaces.",
            "Subscriptions for business analytics, campaigns, and instructor tools as APIs mature.",
            "Partnerships with instructors and brands for curated programs — aligned with user goals.",
        ],
    )

    _bullet_slide(
        prs,
        "Traction & milestones (template)",
        [
            "Replace with your real metrics: MAU/WAU, retention D7/D30, GMV, active businesses.",
            "Near-term: campaign APIs, R2-backed media uploads, analytics depth.",
            "Medium-term: payments integration, moderation tooling, and geographic expansion.",
        ],
    )

    _bullet_slide(
        prs,
        "Competitive positioning",
        [
            "Vs. generic social: growth-oriented categories, CO2-aware framing, and commerce tie-in.",
            "Vs. marketplaces only: community loop that brings repeat visits and organic demand.",
            "Moat over time: graph of progress + instructor trust + business relationships on-platform.",
        ],
    )

    _bullet_slide(
        prs,
        "Team & execution",
        [
            "Highlight founders’ domain expertise in mobile, growth, or retail (edit this slide).",
            "Ship cadence: tested API layer + iterative client releases documented in-repo.",
            "Security posture: JWT sessions, secure storage on device, environment-based configuration.",
        ],
    )

    _bullet_slide(
        prs,
        "The ask",
        [
            "Use of funds: product/engineering, initial growth experiments, compliance, and GTM.",
            "Milestones tied to funding: hire key roles, harden marketplace payments, scale infra.",
            "Contact: add your email, deck link, and data room reference on this slide before sending.",
        ],
    )

    _title_slide(prs, "Thank you", "Growl — grow with your community.")

    prs.save(OUT_PATH)
    print(f"Wrote {OUT_PATH}")


if __name__ == "__main__":
    main()
